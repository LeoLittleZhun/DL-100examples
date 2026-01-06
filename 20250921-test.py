#!/usr/bin/env python
# -*- coding: utf-8 -*-
'''
@File    :   20250921-test.py
@Time    :   2025/09/21 21:56:49
@Author  :   ljc 
@Version :   1.0
@Desc    :   None
'''

# here put the import lib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 创建演示文稿对象
prs = Presentation()

# 设置统一的颜色方案
TITLE_COLOR = RGBColor(0, 102, 204)   # 蓝色
TEXT_COLOR = RGBColor(255, 255, 255)  # 白色
BG_COLOR = RGBColor(15, 35, 70)       # 深蓝背景

# 定义一个通用的页面生成函数
def add_slide(title, content_list, subtitle=None):
    slide_layout = prs.slide_layouts[6]  # 空白页
    slide = prs.slides.add_slide(slide_layout)

    # 背景填充深蓝色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TITLE_COLOR
    title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # 副标题（如果有）
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle
        subtitle_tf.paragraphs[0].font.size = Pt(20)
        subtitle_tf.paragraphs[0].font.color.rgb = TEXT_COLOR

    # 内容
    top = 2 if not subtitle else 2.2
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(9), Inches(5))
    content_tf = content_box.text_frame
    for line in content_list:
        p = content_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0

# === 以下为11页内容 ===

# Page 1: 封面
add_slide(
    "从实践到方法论：复杂项目的可行性分析与风险一体化管控",
    ["高级项目经理 / PMP / CISSP / 注册咨询工程师", "电话：13269855583"],
    "刘佳星 · 8年政企与安全项目交付经验总结"
)

# Page 2: 个人简介
add_slide("个人简介：我的实战履历", [
    "核心画像：善于从0-1交付复杂项目并沉淀体系化方法的高级项目经理",
    "千万级项目：主导4个千万级项目交付，涵盖部委、运营商、大型集团",
    "全周期经验：深度参与项目可研、设计、立项、招投标、实施、验收全流程",
    "跨界能力：兼具技术开发背景（Python/Java/Docker）与项目管理、解决方案能力",
    "今日分享目标：分享我如何系统性地确保项目可行并管控风险"
])

# Page 3: 目录
add_slide("目录与核心逻辑", [
    "总体逻辑：以实践案例为基石，提炼可复制的方法论",
    "1. 核心理念：可行性是前提，风险是伴生 (1 min)",
    "2. 实践框架：我的TEL三维可行性分析模型 (2 min)",
    "3. 风险一体化管控：四大风险域的识别与应对实战 (5 min)",
    "4. 总结：从项目实践到组织资产 (1 min)"
])

# Page 4
add_slide("核心理念：可行性是前提，风险是伴生", [
    "观点：可行性分析不是独立环节，而是持续的风险辨识过程",
    "可行性分析：确定性因素（技术、资源、政策）的深度评估",
    "风险管理：不确定性因素（变化、意外）的动态应对",
    "结合：可行性分析为风险管理划定基线，风险监控反过来验证和修正可行性"
])

# Page 5
add_slide("实践框架：TEL三维可行性分析模型", [
    "技术可行性：能否做得出来？ 技术栈匹配度、系统集成复杂度、性能与安全指标",
    "经济与资源可行性：资源是否足够？预算、人力、时间、供应链稳定性",
    "合规与政策可行性：是否允许做？政策符合性、数据安全、考核指标",
    "案例：安全产品兼容性POC测试；智慧楼宇供应链提前锁定；部委项目天生合规"
])

# Page 6
add_slide("风险一体化管控：四大风险域", [
    "政策与合规风险 -> 外部环境层",
    "客户与需求风险 -> 业务需求层",
    "技术与实施风险 -> 解决方案层",
    "资源与协作风险 -> 执行支撑层"
])

# Page 7
add_slide("风险实战：政策与客户风险", [
    "政策与合规风险应对：主动嵌入，持续监控",
    "案例：集团安全平台，政策嵌入规范，监控机制，柔性设计",
    "客户与需求风险应对：流程管控，范围锚定",
    "案例：能源管理平台，设立CCB，量化impact，需求跟踪矩阵"
])

# Page 8
add_slide("风险实战：技术与资源风险", [
    "技术与实施风险应对：小步验证，透明同步",
    "案例：接口先行、持续集成、日站会",
    "资源与协作风险应对：多源备份，信息牵引",
    "案例：备份机制、资源看板、高层同步"
])

# Page 9
add_slide("方法论沉淀：从个人经验到组织资产", [
    "复盘习惯：经验转化为清单、模板、流程",
    "沉淀资产：《政策合规性检查清单》《变更控制SOP》《供应商评估矩阵》《启动会Checklist》",
    "核心价值：让风险管控不依赖个人，而是成为组织标准动作"
])

# Page 10
add_slide("总结：我的价值", [
    "一套方法：TEL模型 + 四维风险管控",
    "一套工具：清单、模板、流程",
    "一种能力：复杂环境下精准判断、稳健前行",
    "价值：将体系带入团队，提升项目成功率"
])

# Page 11
add_slide("Q&A", ["谢谢聆听！ 欢迎交流！"])

# 保存PPT
prs.save("项目可行性与风险管控.pptx")
print("PPT 文件已生成：项目可行性与风险管控.pptx")